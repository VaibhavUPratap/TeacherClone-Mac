from __future__ import annotations
import os
import uuid
import httpx
import fitz  # PyMuPDF
from datetime import datetime

from config import settings
from .vector_service import vector_service

# Directory where uploaded files are temporarily saved
DOCUMENTS_DIR = os.path.join(os.path.dirname(__file__), "..", "data", "documents")
os.makedirs(DOCUMENTS_DIR, exist_ok=True)

# ── Text chunking helpers ────────────────────────────────────────────────────

CHUNK_WORD_SIZE = 400  # target words per chunk


def _chunk_text(text: str, chunk_size: int = CHUNK_WORD_SIZE) -> list[str]:
    """
    Split *text* into chunks of approximately *chunk_size* words.

    Splits on whitespace; preserves natural word boundaries.
    Returns only non-empty chunks.
    """
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size):
        chunk = " ".join(words[i : i + chunk_size]).strip()
        if chunk:
            chunks.append(chunk)
    return chunks


# ── Embedding helper ─────────────────────────────────────────────────────────

def _get_embedding(text: str) -> list[float]:
    """
    Generate a dense embedding vector for *text* using Ollama's local API.

    Model: nomic-embed-text (configured via OLLAMA_EMBED_MODEL)
    Raises on HTTP errors so callers can surface the problem clearly.
    """
    url = f"{settings.OLLAMA_BASE_URL}/api/embeddings"
    payload = {
        "model": settings.OLLAMA_EMBED_MODEL,
        "prompt": text,
    }
    with httpx.Client(timeout=60.0) as client:
        response = client.post(url, json=payload)
        response.raise_for_status()
    return response.json()["embedding"]


# ── Text extraction ──────────────────────────────────────────────────────────

# Max characters stored in Supabase `resources.content` (prevents huge rows)
_CONTENT_MAX_CHARS = 12_000


def _sanitize_text(text: str) -> str:
    """
    Remove characters that Postgres / Supabase cannot store:
      - Null bytes (\u0000) → '22P05: unsupported Unicode escape sequence'
      - Other non-printable control chars (except \t \n \r)
    Also normalise Windows-style line endings.
    """
    # Strip null bytes and other non-printable control chars
    cleaned = "".join(
        ch for ch in text
        if ch == "\t" or ch == "\n" or ch == "\r" or (ord(ch) >= 32)
    )
    return cleaned.replace("\r\n", "\n").strip()


def _extract_text_from_pdf(file_path: str) -> str:
    """Extract all text from a PDF file using PyMuPDF."""
    doc = fitz.open(file_path)
    pages = [page.get_text() for page in doc]
    doc.close()
    return "\n".join(pages)


def _extract_text_from_pptx(file_path: str) -> str:
    """Extract all text from a PPTX presentation file using python-pptx."""
    from pptx import Presentation
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_runs.append(shape.text.strip())
    return "\n".join(text_runs)


def _extract_text_from_txt(file_path: str) -> str:
    """Read raw text from a .txt file."""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


# ── IngestService ────────────────────────────────────────────────────────────

class IngestService:
    """
    Handles end-to-end document ingestion for the RAG pipeline.

    Flow:
      1. Save the uploaded file to disk.
      2. Extract text (PDF, PPTX, or plain text).
      3. Split text into word-count chunks (~400 words each).
      4. Generate an Ollama embedding for each chunk.
      5. Store chunks + embeddings in ChromaDB via VectorService.
    """

    def process_file(self, filename: str, file_bytes: bytes, subject_id: str = None) -> dict:
        """
        Ingest an uploaded file and populate the vector store.

        Args:
            filename:   Original filename (used to detect type & as metadata).
            file_bytes: Raw file content from the UploadFile.
            subject_id: Optional ID of the subject this file belongs to.

        Returns:
            A dict with file_id, chunk_count, and status.
        """
        # 1. Persist the file in subject-specific subdirectory
        if subject_id:
            subject_dir = os.path.join(DOCUMENTS_DIR, subject_id)
            os.makedirs(subject_dir, exist_ok=True)
            file_path = os.path.join(subject_dir, filename)
        else:
            file_path = os.path.join(DOCUMENTS_DIR, filename)

        with open(file_path, "wb") as f:
            f.write(file_bytes)

        # 2. Extract text
        ext = os.path.splitext(filename)[-1].lower()
        if ext == ".pdf":
            raw_text = _extract_text_from_pdf(file_path)
        elif ext in (".pptx", ".ppt"):
            raw_text = _extract_text_from_pptx(file_path)
        elif ext in (".txt", ".md"):
            raw_text = _extract_text_from_txt(file_path)
        else:
            return {
                "file_id": None,
                "status": "error",
                "detail": f"Unsupported file type: {ext}. Upload PDF, PPTX or TXT.",
            }

        # Sanitize immediately — removes null bytes that break Postgres
        raw_text = _sanitize_text(raw_text)

        if not raw_text.strip():
            return {
                "file_id": None,
                "status": "error",
                "detail": "No text could be extracted from the file (may be a scanned/image-only PDF).",
            }

        # 3. Chunk
        chunks = _chunk_text(raw_text)

        # 4. Embed & 5. Store
        file_id = str(uuid.uuid4())
        ids: list[str] = []
        embeddings: list[list[float]] = []
        metadatas: list[dict] = []

        for i, chunk in enumerate(chunks):
            embedding = _get_embedding(chunk)
            ids.append(f"{file_id}_{i}")
            embeddings.append(embedding)
            # Add subject_id and file_id to metadata so queries can isolate searches by subject
            # and deletions can target specific files.
            metadatas.append({
                "source": filename, 
                "chunk_index": i,
                "subject_id": subject_id or "",
                "file_id": file_id
            })

        vector_service.add_documents(
            chunks=chunks,
            ids=ids,
            embeddings=embeddings,
            metadatas=metadatas,
        )

        from config import supabase
        if supabase is not None:
            try:
                # Insert record into documents table
                supabase.table("documents").insert({
                    "file_id": file_id,
                    "filename": filename,
                    "chunk_count": len(chunks),
                    "timestamp": datetime.now().isoformat(),
                    "status": "completed",
                    "subject_id": subject_id
                }).execute()

                # Also insert into public.resources so it's visible to students in the Library
                ext_type_map = {
                    ".pdf": "Lecture PDF",
                    ".pptx": "Presentation Slides",
                    ".ppt": "Presentation Slides",
                    ".txt": "Study Notes",
                    ".md": "Study Notes"
                }
                res_type = ext_type_map.get(ext, "Study Material")
                
                # Truncate content so large files don't exceed Supabase row limits
                content_preview = raw_text[:_CONTENT_MAX_CHARS]

                supabase.table("resources").insert({
                    "id": file_id,
                    "subject_id": subject_id,
                    "title": filename,
                    "type": res_type,
                    "description": f"Ingested study material. Chunks: {len(chunks)}",
                    "content": content_preview,
                    "created_at": datetime.now().isoformat()
                }).execute()
            except Exception as e:
                print(f"Supabase Ingest Sync Error: {e}")

        return {
            "file_id": file_id,
            "filename": filename,
            "status": "completed"
        }

    @staticmethod
    def get_all_documents():
        """Retrieves all ingested documents from Supabase."""
        from config import supabase
        if supabase is not None:
            try:
                response = supabase.table("documents").select("*").order("timestamp", desc=True).execute()
                return response.data or []
            except Exception as e:
                print(f"Supabase List Documents Error: {e}")
        return []

    def get_ingestion_status(self, file_id: str) -> dict:
        """
        Legacy status endpoint — real status is reflected by chunk_count above.
        """
        return {"file_id": file_id, "status": "completed"}

    def delete_document(self, file_id: str) -> dict:
        """
        Purges a document, its database records (documents and resources),
        its ChromaDB chunks, and its local file.
        """
        from config import supabase
        filename = None
        subject_id = None
        
        # 1. Fetch filename and subject_id from Supabase to resolve physical path and ChromaDB deletion
        if supabase is not None:
            try:
                response = supabase.table("documents").select("filename, subject_id").eq("file_id", file_id).execute()
                if response.data:
                    filename = response.data[0]["filename"]
                    subject_id = response.data[0].get("subject_id")
            except Exception as e:
                print(f"Supabase fetch filename/subject_id error: {e}")
                
        # 2. Delete database records
        if supabase is not None:
            try:
                # Delete from public.documents
                supabase.table("documents").delete().eq("file_id", file_id).execute()
                # Delete from public.resources
                supabase.table("resources").delete().eq("id", file_id).execute()
            except Exception as e:
                print(f"Supabase delete records error: {e}")

        # 3. Purge vector chunks from ChromaDB
        try:
            vector_service.delete_document(file_id=file_id, filename=filename)
        except Exception as e:
            print(f"ChromaDB chunk purge error: {e}")

        # 4. Remove physical file
        if filename:
            if subject_id:
                file_path = os.path.join(DOCUMENTS_DIR, subject_id, filename)
            else:
                file_path = os.path.join(DOCUMENTS_DIR, filename)
                
            # Fallback check at root folder in case it was uploaded before subject folder structuring
            if not os.path.exists(file_path):
                file_path = os.path.join(DOCUMENTS_DIR, filename)

            if os.path.exists(file_path):
                try:
                    os.remove(file_path)
                except Exception as e:
                    print(f"Physical file removal error: {e}")
                    
        return {"status": "success", "detail": f"Document {file_id} deleted successfully."}


ingest_service = IngestService()
